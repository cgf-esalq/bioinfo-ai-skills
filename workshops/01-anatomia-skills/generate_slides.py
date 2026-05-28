#!/usr/bin/env python3
# /// script
# requires-python = ">=3.11"
# dependencies = ["python-pptx"]
# ///
"""
Gera o PPTX do Workshop 1 — Anatomia de uma Skill e o Ecossistema de Tools
bioinfo-ai-skills | CGF-ESALQ
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Paleta de cores ─────────────────────────────────────────────────────────
BG_DARK       = RGBColor(0x0F, 0x0F, 0x23)   # fundo principal
BG_SECTION    = RGBColor(0x1A, 0x1A, 0x2E)   # fundo de seção
ACCENT        = RGBColor(0xA7, 0x8B, 0xFA)   # roxo claro
ACCENT_DARK   = RGBColor(0x7C, 0x3A, 0xED)   # roxo escuro
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xD4, 0xD4, 0xD4)
CODE_BG       = RGBColor(0x1E, 0x1E, 0x2E)
BORDER        = RGBColor(0x31, 0x32, 0x44)
GREEN         = RGBColor(0xA6, 0xE3, 0xA1)

# ── Tamanho widescreen 16:9 ─────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # layout em branco

# ── Helpers ──────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def txbox(slide, text, left, top, width, height,
          font_size=20, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, italic=False, font_name="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = Pt(font_size)
    run.font.bold      = bold
    run.font.color.rgb = color
    run.font.italic    = italic
    run.font.name      = font_name
    return box

def rect(slide, left, top, width, height, fill_color, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width     = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def accent_bar(slide):
    """Barra de acento roxa no topo."""
    rect(slide, 0, 0, W, Inches(0.08), ACCENT_DARK)

def footer_bar(slide, text="🧬 bioinfo-ai-skills | WS1 — Anatomia de Skills"):
    rect(slide, 0, H - Inches(0.35), W, Inches(0.35), BG_SECTION)
    txbox(slide, text,
          Inches(0.3), H - Inches(0.33), W - Inches(0.6), Inches(0.3),
          font_size=10, color=RGBColor(0x6C, 0x70, 0x86), align=PP_ALIGN.CENTER)

def slide_title(slide, title, subtitle=None):
    accent_bar(slide)
    footer_bar(slide)
    txbox(slide, title,
          Inches(0.5), Inches(0.2), W - Inches(1), Inches(0.6),
          font_size=32, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)
    if subtitle:
        txbox(slide, subtitle,
              Inches(0.5), Inches(0.85), W - Inches(1), Inches(0.4),
              font_size=18, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)
    # divider
    rect(slide, Inches(0.5), Inches(1.0), W - Inches(1), Inches(0.025), ACCENT_DARK)

def bullet_items(slide, items, top, left=Inches(0.7), width=None, font_size=18, color=WHITE):
    if width is None:
        width = W - Inches(1.4)
    box = slide.shapes.add_textbox(left, top, width, Inches(5))
    tf  = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = item
        run.font.size      = Pt(font_size)
        run.font.color.rgb = color
        run.font.name      = "Calibri"

def code_box(slide, code, top, left=Inches(0.5), width=None, font_size=13):
    if width is None:
        width = W - Inches(1)
    # estimate height
    lines  = code.count('\n') + 1
    height = Pt(font_size * 1.5 * lines + 20)
    r = rect(slide, left, top, width, height, CODE_BG, BORDER, 0.5)
    box = slide.shapes.add_textbox(left + Inches(0.15), top + Pt(6),
                                   width - Inches(0.3), height - Pt(12))
    tf  = box.text_frame
    tf.word_wrap = False
    p   = tf.paragraphs[0]
    run = p.add_run()
    run.text = code
    run.font.size      = Pt(font_size)
    run.font.color.rgb = GREEN
    run.font.name      = "Courier New"
    return height

def section_slide(slide, number, title):
    bg(slide, BG_SECTION)
    # big number
    txbox(slide, str(number),
          Inches(0.5), Inches(1.5), Inches(2), Inches(2.5),
          font_size=120, bold=True, color=ACCENT_DARK)
    txbox(slide, title,
          Inches(2.8), Inches(2.5), W - Inches(3.3), Inches(1.5),
          font_size=40, bold=True, color=WHITE)
    rect(slide, Inches(2.8), Inches(2.4), Inches(0.06), Inches(1.5), ACCENT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Capa
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
accent_bar(s)

# Gradiente visual: retângulo decorativo lateral
rect(s, W - Inches(4), 0, Inches(4), H, RGBColor(0x1E, 0x1E, 0x3E))
rect(s, W - Inches(4), 0, Inches(0.05), H, ACCENT_DARK)

txbox(s, "🧬 bioinfo-ai-skills",
      Inches(0.6), Inches(1.2), W - Inches(5), Inches(0.6),
      font_size=18, color=ACCENT)
txbox(s, "Anatomia de uma Skill\ne o Ecossistema de Tools",
      Inches(0.6), Inches(1.9), W - Inches(5), Inches(2),
      font_size=38, bold=True, color=WHITE)
txbox(s, "Workshop 1 — CGF-ESALQ",
      Inches(0.6), Inches(4.0), W - Inches(5), Inches(0.5),
      font_size=20, color=ACCENT, bold=True)
txbox(s, "Usando IA de forma responsável em bioinformática genômica",
      Inches(0.6), Inches(4.6), W - Inches(5), Inches(0.5),
      font_size=15, color=LIGHT_GRAY, italic=True)

# Ícone lateral decorativo
txbox(s, "WS1", W - Inches(3.2), Inches(2.8), Inches(2.5), Inches(1.5),
      font_size=72, bold=True, color=ACCENT_DARK, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
slide_title(s, "Agenda")
footer_bar(s)

agenda = [
    ("1", "O que é uma Skill?",                   "25 min"),
    ("2", "Explorando Skills existentes",           "20 min"),
    ("3", "MCP Servers e Tools",                   "20 min"),
    ("☕", "Pausa",                                  " 5 min"),
    ("4", "Hands-on: Buscando dados de Bos taurus","30 min"),
    ("5", "Anatomia detalhada do SKILL.md",         "15 min"),
    ("6", "Discussão: ideias de Skills",            "10 min"),
]

top_start = Inches(1.3)
row_h     = Inches(0.62)
for i, (num, topic, time) in enumerate(agenda):
    top = top_start + i * row_h
    bg_col = RGBColor(0x1E, 0x1E, 0x3E) if i % 2 == 0 else BG_SECTION
    rect(s, Inches(0.4), top, W - Inches(0.8), row_h - Inches(0.04), bg_col)
    # number badge
    rect(s, Inches(0.4), top, Inches(0.7), row_h - Inches(0.04), ACCENT_DARK)
    txbox(s, num, Inches(0.4), top + Pt(2), Inches(0.7), row_h,
          font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s, topic, Inches(1.2), top + Pt(4), Inches(9.5), row_h,
          font_size=16, color=WHITE)
    txbox(s, time, W - Inches(2.0), top + Pt(4), Inches(1.6), row_h,
          font_size=16, color=ACCENT, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Seção 1
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, BG_SECTION)
section_slide(s, 1, "O que é uma Skill?")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — O Problema
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
slide_title(s, "O Problema (e a Solução)")
footer_bar(s)

col_w = Inches(5.6)
# Sem skill
rect(s, Inches(0.4), Inches(1.2), col_w, Inches(5.2), RGBColor(0x2A, 0x1A, 0x1A))
rect(s, Inches(0.4), Inches(1.2), col_w, Inches(0.5), RGBColor(0x7F, 0x1D, 0x1D))
txbox(s, "❌  Sem Skill", Inches(0.5), Inches(1.22), col_w, Inches(0.45),
      font_size=18, bold=True, color=WHITE)
bullet_items(s, [
    "• Você digita o mesmo prompt toda vez",
    "• A IA não sabe suas preferências",
    "• Cada conversa começa do zero",
    "• Resultados inconsistentes",
], top=Inches(1.85), left=Inches(0.6), width=col_w - Inches(0.3), font_size=16)

# Com skill
rect(s, Inches(6.8), Inches(1.2), col_w, Inches(5.2), RGBColor(0x1A, 0x2A, 0x1A))
rect(s, Inches(6.8), Inches(1.2), col_w, Inches(0.5), RGBColor(0x14, 0x53, 0x2D))
txbox(s, "✅  Com Skill", Inches(6.9), Inches(1.22), col_w, Inches(0.45),
      font_size=18, bold=True, color=WHITE)
bullet_items(s, [
    "• Instruções reutilizáveis",
    "• A IA conhece seu contexto",
    "• Comportamento padronizado",
    "• Resultados reprodutíveis",
], top=Inches(1.85), left=Inches(6.9), width=col_w - Inches(0.3), font_size=16)

# VS
txbox(s, "VS", Inches(6.0), Inches(3.4), Inches(0.8), Inches(0.6),
      font_size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Definição
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
slide_title(s, "Definição")
footer_bar(s)

# Quote box
rect(s, Inches(0.4), Inches(1.3), W - Inches(0.8), Inches(2.0),
     RGBColor(0x1E, 0x1E, 0x2E), ACCENT_DARK, 1)
rect(s, Inches(0.4), Inches(1.3), Inches(0.08), Inches(2.0), ACCENT_DARK)
txbox(s, 'Skill = um conjunto de instruções, scripts e recursos que '
         'estendem as capacidades de um agente de IA para tarefas especializadas.',
      Inches(0.7), Inches(1.45), W - Inches(1.4), Inches(1.7),
      font_size=22, color=RGBColor(0xCD, 0xD6, 0xF4))

txbox(s, "💡  Pense numa skill como uma receita detalhada que ensina a IA "
         "a executar uma tarefa específica sempre da mesma forma.",
      Inches(0.5), Inches(3.6), W - Inches(1), Inches(1.5),
      font_size=20, color=LIGHT_GRAY, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Estrutura de uma Skill
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
slide_title(s, "Estrutura de uma Skill")
footer_bar(s)

struct_code = (
    "minha-skill/\n"
    "├── SKILL.md          # Instruções principais (obrigatório)\n"
    "├── scripts/          # Scripts auxiliares\n"
    "│   └── process.R\n"
    "├── resources/        # Dados, templates, referências\n"
    "│   └── reference.md\n"
    "├── examples/         # Exemplos de uso\n"
    "│   └── usage.md\n"
    "└── references/       # Documentação adicional\n"
    "    └── docs.md"
)
code_box(s, struct_code, Inches(0.5), font_size=14)

txbox(s, "⚠️  O único arquivo obrigatório é o SKILL.md — todo o resto é opcional.",
      Inches(0.5), Inches(5.5), W - Inches(1), Inches(0.6),
      font_size=16, color=ACCENT, bold=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — SKILL.md: O Coração
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
slide_title(s, "SKILL.md — O Coração da Skill")
footer_bar(s)

yaml_code = (
    "---\n"
    "name: bovine-variant-filter\n"
    "description: >\n"
    "  Filtra variantes genomicas de Bos taurus (ARS-UCD1.3)\n"
    "  aplicando criterios de qualidade padrao para SNPs e indels.\n"
    "---\n\n"
    "# Instrucoes\n\n"
    "Quando o usuario pedir para filtrar variantes bovinas:\n\n"
    "1. Leia o VCF com VariantAnnotation::readVcf()\n"
    "2. Aplique os filtros: QUAL > 30, DP > 10\n"
    "3. Salve em results/filtered/\n"
    "4. Gere relatorio com estatisticas de filtragem"
)
code_box(s, yaml_code, Inches(1.1), font_size=13)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Por que importa para bioinformática?
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
slide_title(s, "Por que Skills importam para bioinformática?")
footer_bar(s)

reasons = [
    ("🔄", "Reprodutibilidade",  "Mesma análise, mesmos parâmetros, sempre"),
    ("📏", "Padronização",        "Todos no lab seguem o mesmo protocolo"),
    ("⚡", "Eficiência",          "Não precisa re-explicar o contexto ao agente"),
    ("📖", "Documentação",        "A skill É a documentação do método"),
    ("🔗", "Compartilhamento",    "Publique e compartilhe com a comunidade"),
]

for i, (icon, title, desc) in enumerate(reasons):
    top = Inches(1.3) + i * Inches(1.0)
    rect(s, Inches(0.4), top, W - Inches(0.8), Inches(0.85),
         RGBColor(0x1E, 0x1E, 0x3E))
    txbox(s, icon, Inches(0.5), top + Pt(6), Inches(0.7), Inches(0.8),
          font_size=28, align=PP_ALIGN.CENTER)
    txbox(s, title, Inches(1.3), top + Pt(4), Inches(2.5), Inches(0.45),
          font_size=17, bold=True, color=ACCENT)
    txbox(s, desc, Inches(3.9), top + Pt(4), W - Inches(4.5), Inches(0.45),
          font_size=16, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Seção 2
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, BG_SECTION)
section_slide(s, 2, "Explorando Skills Existentes")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Google DeepMind Science Skills
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
slide_title(s, "Google DeepMind Science Skills")
footer_bar(s)

txbox(s, "Conjunto de skills open-source para ciência:",
      Inches(0.5), Inches(1.15), W - Inches(1), Inches(0.45),
      font_size=17, color=LIGHT_GRAY)

code_box(s, "npx -y skills add google-deepmind/science-skills/",
         Inches(1.65), font_size=15)

skills_data = [
    ("science-skills-common", "Utilitários compartilhados"),
    ("string-database",        "Interações proteína-proteína"),
    ("uniprot-database",       "Informações de proteínas"),
    ("unibind-database",       "Fatores de transcrição"),
    ("ucsc-conservation-and-tfbs", "Conservação e TFBSs"),
    ("workflow-skill-creator", "Criar novas skills"),
]

for i, (skill, desc) in enumerate(skills_data):
    top = Inches(2.5) + i * Inches(0.72)
    bg_c = RGBColor(0x1E, 0x1E, 0x3E) if i % 2 == 0 else BG_SECTION
    rect(s, Inches(0.4), top, W - Inches(0.8), Inches(0.66), bg_c)
    txbox(s, skill, Inches(0.6), top + Pt(4), Inches(5.5), Inches(0.55),
          font_size=14, color=GREEN, font_name="Courier New")
    txbox(s, desc, Inches(6.2), top + Pt(4), W - Inches(6.8), Inches(0.55),
          font_size=14, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Exercício: Explore uma skill
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
slide_title(s, "🔍 Exercício: Explore uma Skill")
footer_bar(s)

txbox(s, "5 minutos — Abra o terminal e explore:",
      Inches(0.5), Inches(1.2), W - Inches(1), Inches(0.4),
      font_size=17, color=LIGHT_GRAY)

ex_code = (
    "# Liste todas as skills\n"
    "ls .agents/skills/\n\n"
    "# Escolha uma e explore\n"
    "cat .agents/skills/uniprot-database/SKILL.md\n\n"
    "# Compare com a skill do STRING\n"
    "cat .agents/skills/string-database/SKILL.md"
)
code_box(s, ex_code, Inches(1.7), font_size=14)

txbox(s, "Perguntas para refletir:",
      Inches(0.5), Inches(4.7), W - Inches(1), Inches(0.4),
      font_size=17, bold=True, color=ACCENT)
bullet_items(s, [
    "• Quais seções o SKILL.md tem?",
    "• Existem scripts auxiliares?",
    "• Como as instruções são organizadas?",
], top=Inches(5.15), font_size=15)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Seção 3
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, BG_SECTION)
section_slide(s, 3, "MCP Servers e Tools")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Skills vs. Tools
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
slide_title(s, "Skills vs. Tools — Qual a diferença?")
footer_bar(s)

rows = [
    ("O que é",      "Instruções em texto",           "Código executável"),
    ("Formato",      "Markdown (SKILL.md)",            "JSON schema + servidor MCP"),
    ("Quem executa", "O agente segue as instruções",   "O servidor MCP executa"),
    ("Acesso a dados","Indireto (via agente)",         "Direto (API, banco de dados)"),
    ("Exemplo",      '"Como anotar variantes bovinas"',"gget_blast(sequence)"),
]

header_top = Inches(1.3)
row_h      = Inches(0.88)
# Cabeçalho
for col_i, (label, x, w) in enumerate([
    ("Aspecto", Inches(0.4),  Inches(3.0)),
    ("Skill",   Inches(3.6),  Inches(4.5)),
    ("Tool (MCP)", Inches(8.3), Inches(4.7)),
]):
    rect(s, x, header_top, w, Inches(0.55), ACCENT_DARK)
    txbox(s, label, x + Inches(0.1), header_top + Pt(4), w, Inches(0.5),
          font_size=17, bold=True, color=WHITE)

for i, (aspect, skill_v, tool_v) in enumerate(rows):
    top = header_top + Inches(0.55) + i * row_h
    bg_c = RGBColor(0x1E, 0x1E, 0x3E) if i % 2 == 0 else BG_SECTION
    rect(s, Inches(0.4), top, W - Inches(0.8), row_h - Inches(0.06), bg_c)
    txbox(s, aspect,  Inches(0.5),  top + Pt(4), Inches(2.9), row_h, font_size=14, bold=True, color=ACCENT)
    txbox(s, skill_v, Inches(3.7),  top + Pt(4), Inches(4.3), row_h, font_size=14, color=WHITE)
    txbox(s, tool_v,  Inches(8.4),  top + Pt(4), Inches(4.5), row_h, font_size=14, color=WHITE)

txbox(s, "💡 Skills instruem o agente sobre quando e como usar os tools disponíveis.",
      Inches(0.5), Inches(6.75), W - Inches(1), Inches(0.5),
      font_size=15, color=ACCENT, italic=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Seção 4 (Hands-on)
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, RGBColor(0x15, 0x0A, 0x2E))
section_slide(s, 4, "Hands-on: Buscando Dados de Bos taurus")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15-18 — Exercícios hands-on
# ═══════════════════════════════════════════════════════════════════════════════
exercises = [
    (
        "Exercício 1: Busca de Gene",
        'Peça ao agente:\n"Use o gget para buscar informações do gene DGAT1 em '
        'Bos taurus. Me diga: em qual cromossomo está, quantos transcritos tem, '
        'e qual a sua função principal."',
        ["• O agente escolheu a ferramenta correta?",
         "• A informação retornada é precisa?",
         "• Como você validaria isso?"],
    ),
    (
        "Exercício 2: Busca de Artigos no PubMed",
        '"Busque no PubMed os 5 artigos mais recentes sobre structural variants in '
        'Bos taurus. Liste título, autores e ano."',
        ["• Compare com a busca direta no PubMed no navegador",
         "• Os resultados coincidem?"],
    ),
    (
        "Exercício 3: Interações Proteicas no STRING",
        '"Quais são as 10 principais proteínas que interagem com DGAT1 em bovinos '
        'segundo o STRING database? Qual o score de confiança de cada interação?"',
        ["• O agente usou uma skill, um tool MCP, ou ambos?",
         "• Como você sabe?"],
    ),
    (
        "Exercício 4: Sequência UniProt",
        '"Obtenha a sequência de aminoácidos da proteína DGAT1 bovina do UniProt '
        'em formato FASTA. Salve em data/sample/DGAT1_bos_taurus.fasta."',
        ["• O arquivo foi criado?",
         "• O conteúdo parece correto?"],
    ),
]

for num, (title, prompt, obs) in enumerate(exercises, 1):
    s = add_slide()
    bg(s)
    slide_title(s, title)
    footer_bar(s)

    # prompt box
    rect(s, Inches(0.4), Inches(1.3), W - Inches(0.8), Inches(2.2),
         RGBColor(0x1E, 0x1E, 0x2E), ACCENT_DARK, 0.5)
    rect(s, Inches(0.4), Inches(1.3), Inches(0.08), Inches(2.2), ACCENT)
    txbox(s, prompt, Inches(0.7), Inches(1.4), W - Inches(1.3), Inches(2.0),
          font_size=16, color=RGBColor(0xCD, 0xD6, 0xF4), italic=True)

    txbox(s, "O que observar:", Inches(0.5), Inches(3.8), W - Inches(1), Inches(0.4),
          font_size=17, bold=True, color=ACCENT)
    bullet_items(s, obs, top=Inches(4.3), font_size=16)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Seção 5
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, RGBColor(0x2A, 0x0A, 0x1A))
section_slide(s, 5, "Anatomia Detalhada do SKILL.md")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Frontmatter YAML
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
slide_title(s, "Frontmatter YAML — O que o agente lê primeiro")
footer_bar(s)

yaml_ex = (
    "---\n"
    "name: nome-da-skill\n"
    "description: >\n"
    "  Uma descricao clara e concisa do que a skill faz.\n"
    "  Pode ter multiplas linhas.\n"
    "---"
)
code_box(s, yaml_ex, Inches(1.3), font_size=15)

fields = [
    ("name",        "✅ Obrigatório", "Identificador único da skill"),
    ("description", "✅ Obrigatório", "O que a skill faz — o agente lê primeiro!"),
]
for i, (field, req, desc) in enumerate(fields):
    top = Inches(3.6) + i * Inches(0.9)
    rect(s, Inches(0.4), top, W - Inches(0.8), Inches(0.8),
         RGBColor(0x1E, 0x1E, 0x3E))
    txbox(s, field, Inches(0.6), top + Pt(4), Inches(2.5), Inches(0.7),
          font_size=16, bold=True, color=GREEN, font_name="Courier New")
    txbox(s, req,   Inches(3.2), top + Pt(4), Inches(2.5), Inches(0.7),
          font_size=15, color=ACCENT)
    txbox(s, desc,  Inches(5.8), top + Pt(4), W - Inches(6.4), Inches(0.7),
          font_size=15, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — Boas Práticas
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
slide_title(s, "Boas Práticas para Escrever Skills")
footer_bar(s)

practices = [
    ("1", "Seja específico",       'Use VariantAnnotation::readVcf() > "Leia o VCF"'),
    ("2", "Dê exemplos",           "Inclua inputs e outputs esperados"),
    ("3", "Defina limites",        "O que a skill NÃO deve fazer"),
    ("4", "Use listas numeradas",  "Passos sequenciais claros"),
    ("5", "Inclua validação",      "Como verificar se deu certo"),
    ("6", "Documente parâmetros",  "Genoma, versão, thresholds"),
]

for i, (num, tip, example) in enumerate(practices):
    top = Inches(1.25) + i * Inches(0.96)
    rect(s, Inches(0.4), top, W - Inches(0.8), Inches(0.88),
         RGBColor(0x1E, 0x1E, 0x3E) if i % 2 == 0 else BG_SECTION)
    rect(s, Inches(0.4), top, Inches(0.55), Inches(0.88), ACCENT_DARK)
    txbox(s, num, Inches(0.4), top + Pt(4), Inches(0.55), Inches(0.8),
          font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s, tip, Inches(1.1), top + Pt(4), Inches(3.5), Inches(0.4),
          font_size=16, bold=True, color=ACCENT)
    txbox(s, example, Inches(4.7), top + Pt(4), W - Inches(5.3), Inches(0.8),
          font_size=14, color=LIGHT_GRAY, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — Seção 6
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, BG_SECTION)
section_slide(s, 6, "Discussão: Ideias de Skills")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 23 — Brainstorm
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
slide_title(s, "💬 Brainstorm — Tarefas repetitivas no lab?")
footer_bar(s)

ideas = [
    ("📊", "Filtragem de variantes com os mesmos critérios?"),
    ("🧬", "Anotação funcional sempre no mesmo banco?"),
    ("📈", "Mesmo tipo de gráfico para diferentes datasets?"),
    ("📋", "Relatório com o mesmo formato toda vez?"),
    ("🔄", "Pipeline que roda os mesmos passos?"),
]

for i, (icon, question) in enumerate(ideas):
    top = Inches(1.3) + i * Inches(1.0)
    rect(s, Inches(0.4), top, W - Inches(0.8), Inches(0.88),
         RGBColor(0x1E, 0x1E, 0x3E) if i % 2 == 0 else BG_SECTION)
    txbox(s, icon,     Inches(0.5), top + Pt(6), Inches(0.7), Inches(0.7),
          font_size=28, align=PP_ALIGN.CENTER)
    txbox(s, question, Inches(1.3), top + Pt(8), W - Inches(1.8), Inches(0.7),
          font_size=17, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 24 — Para o próximo workshop
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
slide_title(s, "Para o Próximo Workshop (WS2)")
footer_bar(s)

col_w = Inches(5.8)
# Trazer preparado
rect(s, Inches(0.4), Inches(1.3), col_w, Inches(4.5), RGBColor(0x1E, 0x1E, 0x3E))
rect(s, Inches(0.4), Inches(1.3), col_w, Inches(0.55), ACCENT_DARK)
txbox(s, "📋  Trazer preparado", Inches(0.5), Inches(1.33), col_w, Inches(0.5),
      font_size=17, bold=True, color=WHITE)
bullet_items(s, [
    "• 2-3 ideias de skills para seu projeto",
    "• Um script R que você usa frequentemente",
    "• Acesso ao Antigravity configurado",
], top=Inches(2.0), left=Inches(0.6), width=col_w - Inches(0.3), font_size=16)

# Vamos fazer
rect(s, Inches(7.0), Inches(1.3), col_w, Inches(4.5), RGBColor(0x1E, 0x1E, 0x2E))
rect(s, Inches(7.0), Inches(1.3), col_w, Inches(0.55), RGBColor(0x14, 0x53, 0x2D))
txbox(s, "🔨  Vamos fazer", Inches(7.1), Inches(1.33), col_w, Inches(0.5),
      font_size=17, bold=True, color=WHITE)
bullet_items(s, [
    "• Criar uma skill do zero",
    "• Testar e iterar",
    "• Compartilhar no GitHub",
], top=Inches(2.0), left=Inches(7.1), width=col_w - Inches(0.3), font_size=16)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 25 — Resumo
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
slide_title(s, "Resumo do Workshop")
footer_bar(s)

summary = [
    ("Skill",       "Instruções reutilizáveis para o agente de IA"),
    ("SKILL.md",    "Arquivo principal com frontmatter YAML + instruções em Markdown"),
    ("MCP/Tools",   "Ferramentas que o agente pode executar diretamente"),
    ("Relação",     "Skills instruem o agente sobre como/quando usar os tools"),
]

for i, (concept, desc) in enumerate(summary):
    top = Inches(1.3) + i * Inches(1.2)
    rect(s, Inches(0.4), top, W - Inches(0.8), Inches(1.1),
         RGBColor(0x1E, 0x1E, 0x3E) if i % 2 == 0 else BG_SECTION)
    txbox(s, concept, Inches(0.6), top + Pt(6), Inches(2.5), Inches(0.5),
          font_size=18, bold=True, color=ACCENT)
    txbox(s, desc,    Inches(3.2), top + Pt(6), W - Inches(3.8), Inches(1.0),
          font_size=16, color=WHITE)

rect(s, Inches(0.4), Inches(6.2), W - Inches(0.8), Inches(0.8),
     RGBColor(0x1A, 0x1A, 0x3A))
txbox(s, "🎯 Takeaway: Skills são a ponte entre o conhecimento do especialista (você) "
         "e a capacidade de execução da IA.",
      Inches(0.6), Inches(6.25), W - Inches(1.2), Inches(0.7),
      font_size=16, color=ACCENT, italic=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 26 — Obrigado
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, BG_SECTION)
accent_bar(s)

rect(s, W - Inches(4), 0, Inches(4), H, RGBColor(0x1E, 0x1E, 0x3E))
rect(s, W - Inches(4), 0, Inches(0.05), H, ACCENT_DARK)

txbox(s, "Obrigado! 🧬",
      Inches(0.6), Inches(2.0), W - Inches(5), Inches(1.2),
      font_size=48, bold=True, color=WHITE)

txbox(s, "📂 Material",
      Inches(0.6), Inches(3.5), W - Inches(5), Inches(0.45),
      font_size=18, bold=True, color=ACCENT)
txbox(s, "github.com/cgf-esalq/bioinfo-ai-skills",
      Inches(0.6), Inches(4.0), W - Inches(5), Inches(0.45),
      font_size=16, color=LIGHT_GRAY)

txbox(s, "📧 Dúvidas?",
      Inches(0.6), Inches(4.7), W - Inches(5), Inches(0.45),
      font_size=18, bold=True, color=ACCENT)
txbox(s, "Abra uma issue ou fale com @cavalheiromf",
      Inches(0.6), Inches(5.2), W - Inches(5), Inches(0.45),
      font_size=16, color=LIGHT_GRAY)

txbox(s, "CGF-ESALQ",
      W - Inches(3.5), Inches(3.2), Inches(3), Inches(0.6),
      font_size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Salvar
# ═══════════════════════════════════════════════════════════════════════════════
out = "slides_premium.pptx"
prs.save(out)
print(f"✅ Apresentação salva em: {out}")
print(f"   Total de slides: {len(prs.slides)}")
